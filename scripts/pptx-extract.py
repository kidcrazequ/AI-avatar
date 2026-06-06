#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""python-pptx 纯文本兜底抽取：当 LibreOffice 渲染失败时，至少抽出 PPT 的
文字、表格、演讲者备注，供 opus 结构化。会丢失图片/图表/示意图（已在输出中标注）。

用法：python3 pptx-extract.py <pptx路径>
输出：stdout 打印按页组织的文本（markdown 友好）。
"""
import sys


def main():
    from pptx import Presentation  # 延迟导入，缺包时报清晰错误
    from pptx.util import Emu  # noqa: F401  (确保依赖存在)

    path = sys.argv[1]
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## 幻灯片 {i}")
        had_visual = False
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                rows = list(tbl.rows)
                if rows:
                    cells0 = [c.text.strip() for c in rows[0].cells]
                    out.append("| " + " | ".join(cells0) + " |")
                    out.append("| " + " | ".join("---" for _ in cells0) + " |")
                    for r in rows[1:]:
                        out.append("| " + " | ".join(c.text.strip() for c in r.cells) + " |")
                continue
            if shape.has_chart:
                had_visual = True
                out.append("（此处有图表，纯文本抽取无法还原，需视觉渲染）")
                continue
            if shape.shape_type == 13:  # PICTURE
                had_visual = True
                continue
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    out.append(t)
        if had_visual:
            out.append("（本页含图片/图表等视觉元素，纯文本抽取已略去）")
        # 演讲者备注
        if slide.has_notes_slide:
            note = (slide.notes_slide.notes_text_frame.text or "").strip()
            if note:
                out.append(f"> 备注：{note}")
        out.append("")
    sys.stdout.write("\n".join(out))


if __name__ == "__main__":
    main()
