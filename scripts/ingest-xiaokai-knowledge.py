#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
一次性知识库摄取脚本：把 4 个源目录里的文本类文件机械抽取成 markdown，
落到小凯分身的 knowledge/ 下（按源目录镜像归类）。非文本（CAD/视频/音频/
构建产物/压缩包/图片）跳过并登记到「已跳过」清单（fail loud，不静默丢）。

用法：
  python3 ingest-xiaokai-knowledge.py "<源目录1>" "<源目录2>" ...
输出根：<repo>/avatars/小凯-电气工程师/knowledge/
"""
import os, sys, json, datetime, traceback, subprocess, re

REPO = "/Users/kian/备份/AI/soul"
KN = os.path.join(REPO, "avatars", "小凯-电气工程师", "knowledge")
TODAY = "2026-05-30"

# 源目录 -> knowledge 子目录名
SRC_MAP = {
    "00 国标及技术文献": "国标及技术文献",
    "工作小结": "工作小结",
    "工商业储能": "工商业储能",
    "超充桩": "超充桩",
}

TEXT_EXT = {".pdf",".docx",".doc",".xlsx",".xlsm",".xls",".pptx",".txt",".csv"}
# 明确跳过的非文本/无关类型
SKIP_EXT = {".dwg",".step",".stp",".dwl",".mp4",".avi",".mov",".mp3",".dll",".o",
            ".exe",".dat",".qm",".scu",".release",".debug",".h",".cpp",".temp",
            ".stash",".7z",".zip",".rar",".crdownload",".url",".png",".jpg",
            ".jpeg",".gif",".bmp",".pfd",".0",".release"}

MAX_ROWS = 300   # 单 sheet 抽取行上限
EMPTY_PDF_CHARS = 40  # 文本层少于此值视为扫描件/无文本层

stats = {"converted":0,"skipped":0,"empty_pdf":0,"errors":0}
skipped_list = []   # (path, reason)
error_list = []     # (path, err)
converted_list = [] # (rel_out, source_type)

def fm(source_path, stype, **extra):
    lines = ["---", f"source_path: {source_path}", f"source_type: {stype}", f"ingested: {TODAY}"]
    for k,v in extra.items():
        lines.append(f"{k}: {v}")
    lines.append("---\n")
    return "\n".join(lines)

def write_md(out_path, content):
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    base, ext = os.path.splitext(out_path)
    n = 1
    p = out_path
    while os.path.exists(p):
        p = f"{base}__{n}{ext}"; n += 1
    with open(p, "w", encoding="utf-8") as f:
        f.write(content)
    return p

def conv_pdf(src):
    import fitz
    doc = fitz.open(src)
    npages = doc.page_count
    parts, total = [], 0
    for i, page in enumerate(doc):
        t = page.get_text("text").strip()
        total += len(t)
        if t:
            parts.append(f"\n## 第 {i+1} 页\n\n{t}")
    doc.close()
    body = "\n".join(parts)
    if total < EMPTY_PDF_CHARS:
        stats["empty_pdf"] += 1
        return ("pdf", f"> **注意：该 PDF 文本层为空或极少（疑似扫描件/图片型），需 OCR 后才可入库。当前仅登记来源，未抽取正文。**\n", {"scanned":"true","pages":npages})
    return ("pdf", body, {"pages": npages})

def conv_docx(src):
    import docx
    d = docx.Document(src)
    out = []
    for p in d.paragraphs:
        if p.text.strip():
            style = (p.style.name or "").lower()
            if style.startswith("heading"):
                out.append(f"\n### {p.text.strip()}\n")
            else:
                out.append(p.text.strip())
    for ti, tbl in enumerate(d.tables):
        out.append(f"\n**表 {ti+1}**\n")
        for r in tbl.rows:
            cells = [c.text.strip().replace("\n"," ") for c in r.cells]
            out.append("| " + " | ".join(cells) + " |")
    return ("docx", "\n".join(out), {})

def conv_doc(src):
    # 老 .doc 用 mac textutil 转纯文本
    txt = subprocess.run(["textutil","-convert","txt","-stdout",src],
                         capture_output=True, text=True, timeout=120)
    return ("doc", txt.stdout.strip(), {})

def _rows_to_md(rows):
    out = []
    for i, row in enumerate(rows[:MAX_ROWS]):
        cells = ["" if c is None else str(c).replace("\n"," ").replace("|","/") for c in row]
        if not any(cells): continue
        out.append("| " + " | ".join(cells) + " |")
    if len(rows) > MAX_ROWS:
        out.append(f"\n> （仅抽取前 {MAX_ROWS} 行，共 {len(rows)} 行，余略）")
    return "\n".join(out)

def conv_xlsx(src):
    import openpyxl
    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    nsheets = len(wb.sheetnames)
    out = []
    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows: continue
        out.append(f"\n## Sheet: {ws.title}\n")
        out.append(_rows_to_md(rows))
    wb.close()
    return ("xlsx", "\n".join(out), {"sheets": nsheets})

def conv_xls(src):
    import xlrd
    wb = xlrd.open_workbook(src)
    out = []
    for sh in wb.sheets():
        out.append(f"\n## Sheet: {sh.name}\n")
        rows = [[sh.cell_value(r,c) for c in range(sh.ncols)] for r in range(sh.nrows)]
        out.append(_rows_to_md(rows))
    return ("xls", "\n".join(out), {"sheets": wb.nsheets})

def conv_pptx(src):
    from pptx import Presentation
    prs = Presentation(src)
    out = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shp in slide.shapes:
            if shp.has_text_frame:
                for para in shp.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t: texts.append(t)
            if shp.has_table:
                for r in shp.table.rows:
                    texts.append(" | ".join(c.text.strip() for c in r.cells))
        if texts:
            out.append(f"\n## 幻灯片 {i+1}\n\n" + "\n".join(texts))
    return ("pptx", "\n".join(out), {"slides": len(prs.slides._sldIdLst)})

def conv_text(src):
    with open(src, "r", encoding="utf-8", errors="replace") as f:
        return ("txt", f.read().strip(), {})

CONV = {".pdf":conv_pdf,".docx":conv_docx,".doc":conv_doc,".xlsx":conv_xlsx,
        ".xlsm":conv_xlsx,".xls":conv_xls,".pptx":conv_pptx,".txt":conv_text,".csv":conv_text}

def process_dir(src_root):
    name = os.path.basename(src_root.rstrip("/"))
    sub = SRC_MAP.get(name, name)
    out_root = os.path.join(KN, sub)
    for dirpath, _, files in os.walk(src_root):
        for fn in files:
            src = os.path.join(dirpath, fn)
            ext = os.path.splitext(fn)[1].lower()
            rel = os.path.relpath(src, src_root)
            if ext in SKIP_EXT or fn.startswith("~$") or fn == ".DS_Store":
                stats["skipped"] += 1; skipped_list.append((src, f"非文本/无关类型 {ext or fn}")); continue
            if ext not in CONV:
                stats["skipped"] += 1; skipped_list.append((src, f"未支持类型 {ext}")); continue
            out_path = os.path.join(out_root, os.path.splitext(rel)[0] + ".md")
            # 溯源标签用相对原始目录路径（如 工作小结/xxx.pdf），不写本机绝对前缀，
            # 便于分身包分发后仍可定位原件归属。
            disp = os.path.join(sub, rel)
            try:
                stype, body, extra = CONV[ext](src)
                if not body or not body.strip():
                    body = "> （该文件解析后无可抽取文本内容。）"
                title = os.path.splitext(os.path.basename(rel))[0]
                content = fm(disp, stype, **extra) + f"# {title}\n\n> 来源原件：`{disp}`\n" + body + "\n"
                p = write_md(out_path, content)
                stats["converted"] += 1
                converted_list.append((os.path.relpath(p, KN), stype))
            except Exception as e:
                stats["errors"] += 1
                error_list.append((src, f"{type(e).__name__}: {e}"))

def main():
    srcs = sys.argv[1:]
    for s in srcs:
        if os.path.isdir(s):
            process_dir(s)
        else:
            print("跳过(非目录):", s)
    # 写摄取报告
    rep = os.path.join(KN, "_ingest-report.json")
    with open(rep, "w", encoding="utf-8") as f:
        json.dump({"stats":stats,
                   "converted_count":len(converted_list),
                   "skipped": skipped_list,
                   "errors": error_list,
                   "ingested": TODAY}, f, ensure_ascii=False, indent=1)
    print(json.dumps(stats, ensure_ascii=False))
    print("skipped:", len(skipped_list), "errors:", len(error_list))

if __name__ == "__main__":
    main()
